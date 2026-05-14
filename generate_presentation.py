"""
Generate a minimalist dark-theme PowerPoint presentation for the RAG Mini-Project.
Run with: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0E, 0x0E, 0x0E)   # near-black background
ACCENT    = RGBColor(0xE8, 0xC4, 0x6A)   # warm gold accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x88, 0x88, 0x88)
DIM       = RGBColor(0x2A, 0x2A, 0x2A)   # subtle divider / card bg
CODE_BG   = RGBColor(0x1A, 0x1A, 0x1A)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, font="Helvetica Neue"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return txb


def add_multiline(slide, lines, l, t, w, h,
                  size=18, color=WHITE, spacing=1.15,
                  font="Helvetica Neue", bold=False):
    """lines = list of strings"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.name  = font
        run.font.bold  = bold
    return txb


def accent_bar(slide, t=0.55, height=0.04):
    """Thin horizontal gold line used as a divider / title underline."""
    add_rect(slide, 0.6, t, 1.2, height, ACCENT)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide — large type, minimal."""
    s = blank_slide(prs)
    fill_bg(s)

    # Left accent stripe
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    # Eyebrow label
    add_text(s, "TIES 4911  ·  2026  ·  MiniProject Option 7",
             0.55, 1.6, 10, 0.4, size=13, color=GREY)

    # Main title
    add_text(s, "RAG Q&A System",
             0.55, 2.1, 12, 1.2, size=58, bold=True, color=WHITE)

    # Subtitle
    add_text(s, "Deep Learning Knowledge Base\nwith Retrieval-Augmented Generation",
             0.55, 3.35, 11, 1.0, size=24, color=ACCENT)

    # Bottom author line
    add_text(s, "Amine Bahous",
             0.55, 6.6, 6, 0.5, size=14, color=GREY)


def slide_02_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "01", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Problem", 0.6, 0.35, 8, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    problems = [
        ("Outdated knowledge",   "LLMs have a training cutoff — they cannot answer about recent or domain-specific topics."),
        ("Hallucinations",        "Models confidently generate plausible but incorrect information with no grounding."),
        ("No domain adaptation",  "General-purpose LLMs lack expertise in specialized fields without expensive retraining."),
    ]

    for i, (title, body) in enumerate(problems):
        top = 1.4 + i * 1.7
        add_rect(s, 0.6, top, 11.8, 1.45, DIM)
        add_rect(s, 0.6, top, 0.07, 1.45, ACCENT)
        add_text(s, title, 0.85, top + 0.12, 10, 0.45,
                 size=18, bold=True, color=ACCENT)
        add_text(s, body,  0.85, top + 0.6,  10.8, 0.7,
                 size=15, color=GREY)


def slide_03_solution(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "02", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Solution — RAG", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    add_text(s,
        "Retrieval-Augmented Generation grounds the LLM's answer in a curated "
        "knowledge base — eliminating hallucinations and enabling domain-specific Q&A "
        "without any retraining.",
        0.6, 1.2, 12.0, 0.9, size=17, color=GREY)

    steps = [
        ("INDEX",    "Documents are chunked, embedded\nand stored in a vector database."),
        ("RETRIEVE", "Query is embedded and the top-K\nmost similar chunks are fetched."),
        ("AUGMENT",  "Chunks are injected into the\nLLM prompt as grounding context."),
        ("GENERATE", "LLM answers using ONLY the\nprovided context — no guessing."),
    ]

    for i, (label, desc) in enumerate(steps):
        x = 0.55 + i * 3.18
        add_rect(s, x, 2.3, 2.95, 3.6, DIM)
        add_rect(s, x, 2.3, 2.95, 0.06, ACCENT)

        # Step number
        add_text(s, str(i + 1), x + 0.15, 2.42, 0.5, 0.5,
                 size=11, bold=True, color=ACCENT)

        add_text(s, label, x + 0.18, 2.75, 2.6, 0.5,
                 size=20, bold=True, color=WHITE)
        add_text(s, desc,  x + 0.18, 3.35, 2.62, 1.2,
                 size=14, color=GREY)

        # Arrow between cards
        if i < 3:
            add_text(s, "→", x + 3.0, 3.8, 0.35, 0.4,
                     size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(s, "Original paper: Lewis et al. (2020) — Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks",
             0.6, 6.9, 12, 0.4, size=11, color=GREY, italic=True)


def slide_04_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "03", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Architecture", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    # INDEXING row
    add_text(s, "INDEXING  (offline)", 0.6, 1.25, 5, 0.35, size=12, bold=True, color=ACCENT)

    index_nodes = ["Documents\n(PDF / TXT)", "Chunker\n500 chars", "Embeddings\nMiniLM-L6-v2", "ChromaDB\nvector store"]
    for i, label in enumerate(index_nodes):
        x = 0.6 + i * 3.05
        add_rect(s, x, 1.65, 2.7, 1.1, DIM)
        add_rect(s, x, 1.65, 2.7, 0.05, ACCENT)
        add_text(s, label, x + 0.15, 1.78, 2.4, 0.85, size=14, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, "→", x + 2.72, 2.1, 0.42, 0.4, size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # QUERYING row
    add_text(s, "QUERYING  (online)", 0.6, 3.1, 5, 0.35, size=12, bold=True, color=ACCENT)

    query_nodes = ["User\nQuestion", "Embed\nQuery", "Cosine\nSimilarity", "Top-K\nChunks", "LLM\nAnswer"]
    for i, label in enumerate(query_nodes):
        x = 0.6 + i * 2.44
        add_rect(s, x, 3.5, 2.2, 1.0, DIM)
        add_rect(s, x, 3.5, 2.2, 0.05, ACCENT)
        add_text(s, label, x + 0.1, 3.63, 2.0, 0.75, size=13, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, "→", x + 2.22, 3.9, 0.3, 0.35, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Context injection note
    add_rect(s, 0.6, 4.75, 12.1, 0.9, RGBColor(0x18, 0x18, 0x18))
    add_text(s,
        '  Prompt template:  "Answer ONLY from this context: {retrieved chunks}  ·  Question: {user question}"',
        0.6, 4.88, 12.0, 0.55, size=13, color=GREY, font="Courier New")


def slide_05_techstack(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "04", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Tech Stack", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    rows = [
        ("Document Loading",  "LangChain — PyPDFLoader, TextLoader"),
        ("Text Splitting",    "RecursiveCharacterTextSplitter  (chunk=500, overlap=50)"),
        ("Embeddings",        "sentence-transformers/all-MiniLM-L6-v2  ·  384-dim vectors"),
        ("Vector Store",      "ChromaDB  ·  local persistent storage  ·  cosine similarity"),
        ("LLM Backend",       "Ollama (llama3.2 — local/free)  or  OpenAI GPT  or  retrieval-only"),
        ("Orchestration",     "LangChain  ·  ChatPromptTemplate  ·  LCEL chains"),
        ("Web UI",            "Gradio 4  ·  http://localhost:7860"),
    ]

    for i, (component, detail) in enumerate(rows):
        top = 1.25 + i * 0.82
        add_rect(s, 0.6, top, 12.1, 0.72, DIM)
        add_text(s, component, 0.8, top + 0.12, 2.8, 0.45,
                 size=14, bold=True, color=ACCENT)
        add_text(s, detail, 3.7, top + 0.12, 9.0, 0.45,
                 size=14, color=WHITE, font="Courier New")


def slide_06_knowledge_base(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "05", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Knowledge Base", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    add_text(s, "17 documents indexed  ·  13 course lecture PDFs  +  4 supplementary texts  ·  ChromaDB",
             0.6, 1.15, 12.5, 0.4, size=14, color=GREY)

    # Left column — lecture PDFs
    add_text(s, "COURSE LECTURES  (PDF)", 0.6, 1.65, 6, 0.35, size=11, bold=True, color=ACCENT)
    lectures = [
        ("Lecture 00", "Course Introduction & Overview"),
        ("Lecture 01", "Machine Learning Fundamentals"),
        ("Lecture 02", "Neural Networks & Backpropagation"),
        ("Lecture 03", "Convolutional Neural Networks"),
        ("Lecture 04", "Object Detection — R-CNN family"),
        ("Lecture 05", "Object Detection — YOLO & SSD"),
        ("Lecture 06", "Instance Segmentation"),
        ("Lecture 07", "Recurrent Networks & Sequence Models"),
        ("Lecture 08", "Transformers & Attention"),
        ("Lecture 09", "Large Language Models & RAG"),
        ("Lecture 10", "Generative Models — GAN & VAE"),
        ("Lecture 11", "Anomaly Detection"),
        ("Lecture 12", "Pose Estimation & Tracking"),
    ]
    for i, (lec, topic) in enumerate(lectures):
        top = 2.05 + i * 0.39
        add_rect(s, 0.6, top, 6.1, 0.35, DIM)
        add_text(s, lec,   0.75, top + 0.04, 1.4, 0.28, size=11, bold=True, color=ACCENT)
        add_text(s, topic, 2.2,  top + 0.04, 4.3, 0.28, size=11, color=WHITE)

    # Right column — supplementary texts
    add_text(s, "SUPPLEMENTARY TEXTS", 7.1, 1.65, 6, 0.35, size=11, bold=True, color=ACCENT)
    extras = [
        ("deep_learning_fundamentals.txt",          "ANNs · Backprop · Activations · Optimizers"),
        ("transformers_and_llms.txt",               "BERT · GPT · RAG · Fine-tuning · LoRA"),
        ("computer_vision_and_object_detection.txt","CNNs · YOLO · Segmentation · Tracking"),
        ("anomaly_detection_and_generative_models.txt", "Autoencoders · VAE · GAN · Diffusion"),
    ]
    for i, (fname, topics) in enumerate(extras):
        top = 2.05 + i * 1.3
        add_rect(s, 7.1, top, 5.9, 1.15, DIM)
        add_rect(s, 7.1, top, 0.06, 1.15, ACCENT)
        add_text(s, fname,  7.3, top + 0.1,  5.5, 0.38, size=11, bold=True, color=WHITE, font="Courier New")
        add_text(s, topics, 7.3, top + 0.55, 5.5, 0.45, size=12, color=GREY)


def slide_07_demo(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "06", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Demo", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    add_text(s, "Web interface at  http://127.0.0.1:7860  (Gradio)",
             0.6, 1.15, 10, 0.4, size=14, color=GREY)

    questions = [
        "What is self-attention and how does it work?",
        "How does YOLO detect objects in an image?",
        "What is the difference between VAE and GAN?",
        "How do autoencoders detect anomalies?",
        "What is LoRA fine-tuning and why is it efficient?",
        "What metrics are used to evaluate object detection?",
    ]

    add_text(s, "Example questions asked:", 0.6, 1.65, 6, 0.4,
             size=15, bold=True, color=ACCENT)

    for i, q in enumerate(questions):
        top = 2.1 + i * 0.72
        add_rect(s, 0.6, top, 8.5, 0.58, DIM)
        add_text(s, f"  \"{q}\"", 0.6, top + 0.08, 8.5, 0.42,
                 size=13, color=WHITE, font="Courier New")

    # Side note
    add_rect(s, 9.4, 1.65, 3.5, 5.1, RGBColor(0x14, 0x14, 0x14))
    add_rect(s, 9.4, 1.65, 3.5, 0.05, ACCENT)
    add_multiline(s, [
        "UI Features",
        "",
        "→  Top-K slider (1–8)",
        "→  Source citations",
        "→  Relevance scores",
        "→  Full context toggle",
        "→  Example questions",
        "",
        "LLM backend swappable",
        "via .env — no code change",
    ], 9.55, 1.75, 3.2, 4.8, size=13, color=GREY)


def slide_08_results(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "07", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Results", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    metrics = [
        ("73+",        "chunks indexed\nacross 17 documents"),
        ("384",        "embedding dimensions\n(MiniLM-L6-v2)"),
        ("~18s",       "total ingestion time\nincl. embedding"),
        ("< 1s",       "query response time\n(retrieval-only mode)"),
    ]

    for i, (value, label) in enumerate(metrics):
        x = 0.6 + i * 3.18
        add_rect(s, x, 1.3, 2.95, 2.0, DIM)
        add_rect(s, x, 1.3, 2.95, 0.05, ACCENT)
        add_text(s, value, x + 0.15, 1.5, 2.65, 0.85,
                 size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, x + 0.15, 2.3, 2.65, 0.75,
                 size=13, color=GREY, align=PP_ALIGN.CENTER)

    benefits = [
        ("No hallucinations",    "Answers are grounded exclusively in the provided documents."),
        ("No retraining needed", "Add new PDFs and re-index in seconds — no GPU required."),
        ("Transparent",          "Every answer shows which source file and chunk it came from."),
        ("Flexible LLM backend", "Swap between Ollama, OpenAI, or retrieval-only via one config line."),
    ]

    add_text(s, "Key Benefits", 0.6, 3.55, 4, 0.4, size=16, bold=True, color=WHITE)

    for i, (title, body) in enumerate(benefits):
        top = 4.05 + i * 0.72
        add_rect(s, 0.6, top, 12.1, 0.6, DIM)
        add_rect(s, 0.6, top, 0.06, 0.6, ACCENT)
        add_text(s, title, 0.82, top + 0.08, 2.8, 0.42, size=13, bold=True, color=WHITE)
        add_text(s, body,  3.7,  top + 0.08, 9.0, 0.42, size=13, color=GREY)


def slide_09_conclusion(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 0.06, 7.5, ACCENT)

    add_text(s, "08", 11.8, 0.3, 1.2, 0.7, size=52, bold=True, color=DIM)
    add_text(s, "Conclusion", 0.6, 0.35, 9, 0.65, size=40, bold=True, color=WHITE)
    accent_bar(s, t=1.05)

    add_text(s,
        "RAG is a powerful, practical technique that bridges the gap between\n"
        "general-purpose LLMs and domain-specific knowledge — without fine-tuning.",
        0.6, 1.2, 12.0, 1.0, size=18, color=GREY)

    takeaways = [
        "Retrieval decouples knowledge from model weights — update the DB, not the model.",
        "sentence-transformers provide strong semantic search with minimal compute.",
        "ChromaDB enables a fully local, zero-cost vector store.",
        "The same pipeline scales to thousands of PDFs with no architectural changes.",
        "Switching LLM backends (Ollama ↔ OpenAI) requires changing one config line.",
    ]

    for i, point in enumerate(takeaways):
        top = 2.35 + i * 0.82
        add_rect(s, 0.6, top, 12.1, 0.68, DIM)
        add_rect(s, 0.6, top, 0.06, 0.68, ACCENT)
        add_text(s, f"  {point}", 0.7, top + 0.11, 11.8, 0.46, size=14, color=WHITE)

    add_text(s, "Lewis et al. (2020) · Vaswani et al. (2017) · LangChain · ChromaDB · Ollama",
             0.6, 7.0, 12, 0.35, size=11, color=GREY, italic=True)


def slide_10_closing(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Full-width gold bar at bottom
    add_rect(s, 0, 6.9, 13.33, 0.6, ACCENT)

    add_text(s, "Thank you", 1.0, 1.8, 11, 1.6,
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "Deep Learning Knowledge Base — RAG Q&A System",
             1.0, 3.5, 11, 0.6, size=22, color=GREY, align=PP_ALIGN.CENTER)

    add_text(s, "TIES 4911  ·  MiniProject Option 7  ·  2026",
             1.0, 4.2, 11, 0.5, size=16, color=GREY, align=PP_ALIGN.CENTER)

    add_text(s, "github  ·  Amine Bahous",
             1.0, 6.95, 11, 0.4, size=13,
             color=RGBColor(0x0E, 0x0E, 0x0E),
             bold=True, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_architecture(prs)
    slide_05_techstack(prs)
    slide_06_knowledge_base(prs)
    slide_07_demo(prs)
    slide_08_results(prs)
    slide_09_conclusion(prs)
    slide_10_closing(prs)

    out = "RAG_MiniProject_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
